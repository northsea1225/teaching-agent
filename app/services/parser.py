from __future__ import annotations

from pathlib import Path

import fitz
from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape

from app.models import ParsedAsset, ResourceType
from app.services.audio import get_audio_duration_seconds


EXTENSION_TO_RESOURCE_TYPE = {
    ".pdf": ResourceType.PDF,
    ".docx": ResourceType.DOCX,
    ".pptx": ResourceType.PPTX,
    ".png": ResourceType.IMAGE,
    ".jpg": ResourceType.IMAGE,
    ".jpeg": ResourceType.IMAGE,
    ".bmp": ResourceType.IMAGE,
    ".gif": ResourceType.IMAGE,
    ".webp": ResourceType.IMAGE,
    ".wav": ResourceType.AUDIO,
    ".mp3": ResourceType.AUDIO,
    ".m4a": ResourceType.AUDIO,
    ".aac": ResourceType.AUDIO,
    ".ogg": ResourceType.AUDIO,
    ".mp4": ResourceType.VIDEO,
    ".mov": ResourceType.VIDEO,
    ".avi": ResourceType.VIDEO,
    ".txt": ResourceType.TEXT,
    ".md": ResourceType.TEXT,
}


def detect_resource_type(filename: str, content_type: str | None = None) -> ResourceType | None:
    suffix = Path(filename).suffix.lower()
    if suffix in EXTENSION_TO_RESOURCE_TYPE:
        return EXTENSION_TO_RESOURCE_TYPE[suffix]

    if content_type:
        if content_type.startswith("image/"):
            return ResourceType.IMAGE
        if content_type.startswith("audio/"):
            return ResourceType.AUDIO
        if content_type.startswith("video/"):
            return ResourceType.VIDEO
        if content_type == "application/pdf":
            return ResourceType.PDF
        if content_type.startswith("text/"):
            return ResourceType.TEXT
    return None


def _preview(text: str, limit: int = 180) -> str:
    normalized = " ".join(text.split())
    return normalized[:limit]


def parse_pdf(path: Path) -> ParsedAsset:
    doc = fitz.open(path)
    page_texts: list[str] = []
    for page in doc:
        page_texts.append(page.get_text("text"))
    extracted_text = "\n".join(page_texts).strip()
    return ParsedAsset(
        resource_type=ResourceType.PDF,
        source_path=str(path),
        extracted_text=extracted_text,
        text_preview=_preview(extracted_text) or f"PDF with {doc.page_count} page(s)",
        page_count=doc.page_count,
        metadata={
            "page_count": doc.page_count,
            "filename": path.name,
        },
    )


def parse_image(path: Path) -> ParsedAsset:
    with Image.open(path) as image:
        width, height = image.size
        image_format = image.format or path.suffix.lstrip(".").upper()
        mode = image.mode
    preview = f"Image {path.name} ({width}x{height}, {image_format}, mode={mode})"
    return ParsedAsset(
        resource_type=ResourceType.IMAGE,
        source_path=str(path),
        text_preview=preview,
        width=width,
        height=height,
        metadata={
            "filename": path.name,
            "format": image_format,
            "mode": mode,
            "width": width,
            "height": height,
        },
        warnings=["OCR not implemented yet; only image metadata is extracted."],
    )


def parse_audio(path: Path) -> ParsedAsset:
    duration_seconds = get_audio_duration_seconds(path)
    preview = f"Audio file {path.name}"
    if duration_seconds is not None:
        preview = f"{preview} ({duration_seconds} seconds)"
    warnings: list[str] = []
    if duration_seconds is None:
        warnings.append("Could not determine audio duration.")
    warnings.append("Transcription is not implemented yet; only audio metadata is extracted.")
    return ParsedAsset(
        resource_type=ResourceType.AUDIO,
        source_path=str(path),
        text_preview=preview,
        duration_seconds=duration_seconds,
        metadata={
            "filename": path.name,
            "duration_seconds": duration_seconds,
        },
        warnings=warnings,
    )


def parse_docx(path: Path) -> ParsedAsset:
    document = Document(path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_cells: list[str] = []
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    table_cells.append(text)

    extracted_parts = paragraphs + table_cells
    extracted_text = "\n".join(extracted_parts).strip()
    return ParsedAsset(
        resource_type=ResourceType.DOCX,
        source_path=str(path),
        extracted_text=extracted_text,
        text_preview=_preview(extracted_text) or f"DOCX file {path.name}",
        metadata={
            "filename": path.name,
            "paragraph_count": len(paragraphs),
            "table_cell_count": len(table_cells),
        },
    )


def _extract_shape_text(shape: BaseShape) -> list[str]:
    texts: list[str] = []
    if isinstance(shape, GroupShape):
        for inner_shape in shape.shapes:
            texts.extend(_extract_shape_text(inner_shape))
        return texts

    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    texts.append(text)
    return texts


def parse_pptx(path: Path) -> ParsedAsset:
    presentation = Presentation(path)
    slide_texts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_extract_shape_text(shape))
        if parts:
            slide_texts.append(f"[Slide {slide_index}]\n" + "\n".join(parts))

    extracted_text = "\n\n".join(slide_texts).strip()
    return ParsedAsset(
        resource_type=ResourceType.PPTX,
        source_path=str(path),
        extracted_text=extracted_text,
        text_preview=_preview(extracted_text) or f"PPTX file {path.name}",
        page_count=len(presentation.slides),
        metadata={
            "filename": path.name,
            "slide_count": len(presentation.slides),
        },
    )


def parse_text(path: Path) -> ParsedAsset:
    extracted_text = path.read_text(encoding="utf-8", errors="ignore")
    return ParsedAsset(
        resource_type=ResourceType.TEXT,
        source_path=str(path),
        extracted_text=extracted_text,
        text_preview=_preview(extracted_text) or f"Text file {path.name}",
        metadata={"filename": path.name},
    )


def parse_file(path: Path, resource_type: ResourceType | None = None) -> ParsedAsset:
    file_type = resource_type or detect_resource_type(path.name)
    if file_type is None:
        raise ValueError(f"Unsupported file type: {path.suffix or path.name}")

    if file_type == ResourceType.PDF:
        return parse_pdf(path)
    if file_type == ResourceType.DOCX:
        return parse_docx(path)
    if file_type == ResourceType.PPTX:
        return parse_pptx(path)
    if file_type == ResourceType.IMAGE:
        return parse_image(path)
    if file_type == ResourceType.AUDIO:
        return parse_audio(path)
    if file_type == ResourceType.TEXT:
        return parse_text(path)

    raise ValueError(f"Parsing not implemented yet for resource type: {file_type.value}")
