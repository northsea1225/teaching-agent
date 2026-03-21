from __future__ import annotations

from datetime import datetime, timezone

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.models import Citation, ExportArtifact, ResourceType, SessionStage, SessionState, SvgBlockSpec, SvgSlideSpec, TeachingSpec
from app.services.planner import generate_slide_plan_for_session
from app.services.svg import generate_svg_deck_for_session
from app.services.storage import attach_export_to_session
from app.utils.paths import build_export_path, ensure_project_directories


PPT_CITATION_THEME_STYLES: dict[str, dict[str, str | float]] = {
    "academy": {
        "heading_fill": "#16324f",
        "heading_text": "#ffffff",
        "heading_fill_transparency": 0.04,
        "chip_fill": "#d7e6f5",
        "chip_text": "#16324f",
        "chip_line": "#16324f",
        "chip_fill_transparency": 0.18,
        "chip_line_transparency": 0.48,
    },
    "studio": {
        "heading_fill": "#8a2c0d",
        "heading_text": "#fffaf7",
        "heading_fill_transparency": 0.02,
        "chip_fill": "#fde6d8",
        "chip_text": "#6b240c",
        "chip_line": "#8a2c0d",
        "chip_fill_transparency": 0.08,
        "chip_line_transparency": 0.2,
    },
    "field": {
        "heading_fill": "#166534",
        "heading_text": "#f8fff9",
        "heading_fill_transparency": 0.03,
        "chip_fill": "#d7f0df",
        "chip_text": "#1c4b32",
        "chip_line": "#166534",
        "chip_fill_transparency": 0.1,
        "chip_line_transparency": 0.24,
    },
    "briefing": {
        "heading_fill": "#0f172a",
        "heading_text": "#f8fafc",
        "heading_fill_transparency": 0.0,
        "chip_fill": "#e2e8f0",
        "chip_text": "#172033",
        "chip_line": "#334155",
        "chip_fill_transparency": 0.14,
        "chip_line_transparency": 0.34,
    },
}


def utc_now() -> datetime:
    return datetime.now(timezone.utc)


def _add_bullet_list(document: Document, items: list[str]) -> None:
    for item in items:
        document.add_paragraph(item, style="List Bullet")


def _overview_lines(spec: TeachingSpec) -> list[str]:
    lines: list[str] = []
    if spec.education_stage:
        lines.append(f"学段：{spec.education_stage}")
    if spec.grade_level:
        lines.append(f"年级：{spec.grade_level}")
    if spec.subject:
        lines.append(f"学科：{spec.subject}")
    if spec.lesson_title:
        lines.append(f"课题：{spec.lesson_title}")
    if spec.class_duration_minutes:
        lines.append(f"课时长度：{spec.class_duration_minutes} 分钟")
    if spec.lesson_count:
        lines.append(f"课时数：{spec.lesson_count}")
    if spec.style_preferences:
        lines.append(f"风格偏好：{', '.join(spec.style_preferences)}")
    if spec.interaction_preferences:
        lines.append("互动偏好：" + ", ".join(mode.value for mode in spec.interaction_preferences))
    return lines


def _build_summary(spec: TeachingSpec, slide_count: int) -> str:
    lesson_title = spec.lesson_title or "未命名课题"
    subject = spec.subject or "general"
    stage = spec.education_stage or "general-stage"
    return f"{lesson_title} DOCX 教案草稿，包含 {slide_count} 页逐页策划，适用于 {stage} {subject} 场景。"


def _build_pptx_summary(spec: TeachingSpec, slide_count: int) -> str:
    lesson_title = spec.lesson_title or "未命名课题"
    subject = spec.subject or "general"
    stage = spec.education_stage or "general-stage"
    return f"{lesson_title} PPTX 课件草稿，包含 {slide_count} 页基于 SVG 中间层渲染的页面，适用于 {stage} {subject} 场景。"


def _ensure_slide_plan(
    session: SessionState,
    store_namespace: str | None,
    top_k: int,
) -> SessionState:
    if session.teaching_spec is None:
        raise ValueError("Session has no teaching spec")
    if session.slide_plan is None:
        session = generate_slide_plan_for_session(
            session,
            store_namespace=store_namespace,
            top_k=top_k,
        )
    assert session.slide_plan is not None
    return session


def _ensure_svg_deck(
    session: SessionState,
    store_namespace: str | None,
    top_k: int,
    theme_id: str | None = None,
    font_preset: str | None = None,
) -> SessionState:
    session = _ensure_slide_plan(session, store_namespace, top_k)
    if session.svg_deck is None:
        session = generate_svg_deck_for_session(
            session,
            store_namespace=store_namespace,
            top_k=top_k,
            theme_id=theme_id,
            font_preset=font_preset,
        )
    assert session.svg_deck is not None
    return session


def _finalize_export(
    session: SessionState,
    artifact: ExportArtifact,
) -> tuple[SessionState, ExportArtifact]:
    session = attach_export_to_session(session, artifact)
    session.stage = SessionStage.EXPORT
    session.last_summary = artifact.summary
    session.updated_at = utc_now()
    return session, artifact


def export_docx_for_session(
    session: SessionState,
    store_namespace: str | None = None,
    top_k: int = 5,
) -> tuple[SessionState, ExportArtifact]:
    session = _ensure_slide_plan(session, store_namespace, top_k)

    ensure_project_directories()
    lesson_title = session.teaching_spec.lesson_title or "lesson"
    export_path = build_export_path(session.session_id, f"{lesson_title}_lesson_plan", "docx")

    document = Document()
    document.core_properties.title = f"{lesson_title} lesson package"
    document.core_properties.subject = session.teaching_spec.subject or "teaching-agent"
    document.core_properties.comments = "Generated by Teaching Agent"

    document.add_heading(lesson_title, level=0)

    overview = _overview_lines(session.teaching_spec)
    if overview:
        document.add_paragraph(" | ".join(overview))

    document.add_heading("约束确认", level=1)
    document.add_paragraph(session.planning_confirmation.summary or "当前还没有约束确认结果。")
    if session.planning_confirmation.items:
        _add_bullet_list(
            document,
            [
                f"{item.label}：{item.detail or '待补充'}"
                for item in session.planning_confirmation.items
            ],
        )

    if session.quality_report is not None:
        document.add_heading("质量检查", level=1)
        document.add_paragraph(
            f"状态：{session.quality_report.status} | 分数：{session.quality_report.score}"
        )
        if session.quality_report.summary:
            document.add_paragraph(session.quality_report.summary)
        if session.quality_report.issues:
            _add_bullet_list(
                document,
                [
                    f"{issue.severity.upper()} · {issue.message}"
                    + (f"（第 {issue.slide_number} 页）" if issue.slide_number else "")
                    for issue in session.quality_report.issues[:8]
                ],
            )

    if session.outline is not None:
        document.add_heading("课程概览", level=1)
        if session.outline.summary:
            document.add_paragraph(session.outline.summary)
        for section in session.outline.sections:
            document.add_heading(section.title, level=2)
            document.add_paragraph(f"目标：{section.goal}")
            if section.bullet_points:
                _add_bullet_list(document, section.bullet_points)

    document.add_heading("逐页策划", level=1)
    for slide in session.slide_plan.slides:
        document.add_heading(f"第 {slide.slide_number} 页：{slide.title}", level=2)
        document.add_paragraph(f"页面类型：{slide.slide_type.value}")
        document.add_paragraph(f"页面目标：{slide.goal}")
        document.add_paragraph(f"互动方式：{slide.interaction_mode.value}")
        if slide.layout_hint:
            document.add_paragraph(f"版式建议：{slide.layout_hint}")

        if slide.key_points:
            document.add_paragraph("关键要点：")
            _add_bullet_list(document, slide.key_points)
        if slide.visual_brief:
            document.add_paragraph("视觉提示：")
            _add_bullet_list(document, slide.visual_brief)
        if slide.speaker_notes:
            document.add_paragraph("讲解备注：")
            _add_bullet_list(document, slide.speaker_notes)
        if slide.citations:
            document.add_paragraph("引用资料：")
            _add_bullet_list(
                document,
                [
                    citation.note or citation.page_label or citation.asset_id
                    for citation in slide.citations
                ],
            )

    if session.uploaded_files:
        document.add_heading("当前会话上传资料", level=1)
        _add_bullet_list(
            document,
            [
                f"{session_file.filename} ({session_file.resource_type.value})"
                for session_file in session.uploaded_files
            ],
        )

    if session.retrieval_hits:
        document.add_heading("知识库参考片段", level=1)
        _add_bullet_list(
            document,
            [" ".join(hit.content.split())[:120] for hit in session.retrieval_hits[:5]],
        )

    document.save(export_path)

    artifact = ExportArtifact(
        filename=export_path.name,
        resource_type=ResourceType.DOCX,
        path=str(export_path),
        summary=_build_summary(session.teaching_spec, len(session.slide_plan.slides)),
    )
    return _finalize_export(session, artifact)


def _rgb(hex_color: str | RGBColor | None, fallback: str = "#16324f") -> RGBColor:
    if isinstance(hex_color, RGBColor):
        return hex_color
    color = (hex_color or fallback).strip().lstrip("#")
    if len(color) == 3:
        color = "".join(ch * 2 for ch in color)
    if len(color) != 6:
        color = fallback.lstrip("#")
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _scale_x(presentation: Presentation, value: int, base_width: int) -> Emu:
    return Emu(int((value / base_width) * presentation.slide_width))


def _scale_y(presentation: Presentation, value: int, base_height: int) -> Emu:
    return Emu(int((value / base_height) * presentation.slide_height))


def _shape_type(variant: str) -> MSO_AUTO_SHAPE_TYPE:
    if variant in {
        "hero",
        "hero-bar",
        "soft-card",
        "outline-card",
        "card",
        "spotlight",
        "glass-card",
        "editorial-panel",
        "chip",
    }:
        return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    if variant == "strip":
        return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE


def _shape_fill_transparency(variant: str) -> float:
    mapping = {
        "glass-card": 0.18,
        "editorial-panel": 0.04,
        "outline-card": 0.82,
        "chip": 0.03,
        "strip": 0.06,
    }
    return mapping.get(variant, 0.0)


def _shape_line_transparency(variant: str) -> float:
    mapping = {
        "outline-card": 0.0,
        "glass-card": 0.55,
        "editorial-panel": 0.28,
        "chip": 0.45,
        "strip": 0.95,
        "hero-bar": 0.22,
        "spotlight": 0.12,
    }
    return mapping.get(variant, 0.18)


def _add_overlay_shape(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    *,
    x: int,
    y: int,
    width: int,
    height: int,
    fill_color: str | RGBColor,
    transparency: float,
    shape_type: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
) -> None:
    overlay = slide.shapes.add_shape(
        shape_type,
        _scale_x(presentation, x, slide_spec.width),
        _scale_y(presentation, y, slide_spec.height),
        _scale_x(presentation, width, slide_spec.width),
        _scale_y(presentation, height, slide_spec.height),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = _rgb(fill_color)
    overlay.fill.transparency = transparency
    overlay.line.fill.background()


def _add_block_underlay(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    block: SvgBlockSpec,
) -> None:
    if block.shape_variant != "spotlight":
        return

    underlay_x = max(block.x - 12, 0)
    underlay_y = max(block.y - 10, 0)
    underlay_width = min(block.width + 18, slide_spec.width - underlay_x)
    underlay_height = min(block.height + 18, slide_spec.height - underlay_y)
    _add_overlay_shape(
        presentation,
        slide,
        slide_spec,
        x=underlay_x,
        y=underlay_y,
        width=underlay_width,
        height=underlay_height,
        fill_color=block.accent_color or slide_spec.accent_color,
        transparency=0.84,
    )


def _add_block_accents(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    block: SvgBlockSpec,
) -> None:
    accent = block.accent_color or slide_spec.accent_color
    if block.shape_variant == "editorial-panel":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=block.x,
            y=block.y,
            width=block.width,
            height=14,
            fill_color=accent,
            transparency=0.0,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
    elif block.shape_variant == "hero-bar":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=block.x,
            y=block.y,
            width=18,
            height=block.height,
            fill_color=accent,
            transparency=0.0,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
    elif block.shape_variant == "glass-card":
        inset_width = max(block.width - 36, 24)
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=block.x + 18,
            y=block.y + 16,
            width=inset_width,
            height=12,
            fill_color="#ffffff",
            transparency=0.78,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )


def _apply_block_shape_style(shape, block: SvgBlockSpec, slide_spec: SvgSlideSpec) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(block.background_fill, "#ffffff")
    shape.fill.transparency = _shape_fill_transparency(block.shape_variant)
    if block.shape_variant == "strip":
        shape.line.fill.background()
        return

    shape.line.color.rgb = _rgb(block.stroke_color or block.accent_color or slide_spec.accent_color, "#16324f")
    shape.line.transparency = _shape_line_transparency(block.shape_variant)


def _apply_block_text_style(frame, block: SvgBlockSpec, slide_spec: SvgSlideSpec) -> None:
    frame.word_wrap = True
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_top = Pt(8)
    frame.margin_bottom = Pt(8)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    title_alignment = PP_ALIGN.LEFT
    body_alignment = PP_ALIGN.LEFT
    title_size = block.title_size
    body_size = block.body_size
    title_space_after = Pt(6)
    body_space_after = Pt(5)
    title_line_spacing = 1.0
    body_line_spacing = 1.14
    title_color = _rgb(block.accent_color or slide_spec.accent_color, "#16324f")
    body_color = _rgb(block.text_color, "#213446")

    if block.role == "hero":
        frame.margin_left = Pt(18)
        frame.margin_right = Pt(18)
        frame.margin_top = Pt(18)
        frame.margin_bottom = Pt(16)
        title_size = max(block.title_size + 10, 30)
        body_size = max(block.body_size + 1, 18)
        title_space_after = Pt(12)
        body_space_after = Pt(7)
        body_line_spacing = 1.2
        title_color = _rgb(slide_spec.text_color, "#17202a")
    elif block.role in {"snapshot", "rhythm", "source-lens"}:
        frame.margin_left = Pt(12)
        frame.margin_right = Pt(12)
        frame.margin_top = Pt(10)
        frame.margin_bottom = Pt(10)
        title_size = max(block.title_size - 3, 14)
        body_size = max(block.body_size - 2, 12)
        title_space_after = Pt(7)
        body_space_after = Pt(5)
    elif block.role in {"footer-strip", "takeaway-strip", "evidence-strip"}:
        title_size = max(block.title_size - 5, 12)
        body_size = max(block.body_size - 4, 11)
        title_space_after = Pt(3)
        body_space_after = Pt(2)
        body_line_spacing = 1.08
    elif block.role in {"launch-chip", "compare-chip", "assignment-tag"}:
        title_alignment = PP_ALIGN.CENTER
        body_alignment = PP_ALIGN.CENTER
        title_size = max(block.title_size - 2, 11)
        body_size = max(block.body_size - 2, 11)

    if block.shape_variant == "chip":
        frame.margin_left = Pt(6)
        frame.margin_right = Pt(6)
        frame.margin_top = Pt(4)
        frame.margin_bottom = Pt(4)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_alignment = PP_ALIGN.CENTER
        title_size = max(block.title_size - 1, 12)
        body_alignment = PP_ALIGN.CENTER
        title_space_after = Pt(1)
        body_space_after = Pt(1)
        body_line_spacing = 1.0
    elif block.shape_variant == "hero-bar":
        frame.margin_left = Pt(18)
        frame.margin_right = Pt(10)
        frame.margin_top = Pt(6)
        frame.margin_bottom = Pt(6)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_size = min(max(title_size - 5, 12), 16)
        title_space_after = Pt(0)
        body_space_after = Pt(0)
        body_line_spacing = 1.0
    elif block.shape_variant == "strip":
        frame.margin_left = Pt(8)
        frame.margin_right = Pt(8)
        frame.margin_top = Pt(4)
        frame.margin_bottom = Pt(4)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_size = max(title_size - 3, 12)
        body_size = max(body_size - 2, 11)
        title_space_after = Pt(2)
        body_space_after = Pt(2)
        body_line_spacing = 1.02
    elif block.shape_variant == "spotlight":
        title_size = block.title_size + 2
        title_space_after = Pt(10)
        body_space_after = Pt(6)
        body_line_spacing = 1.18
    elif block.shape_variant == "glass-card":
        body_space_after = Pt(6)
        body_line_spacing = 1.18
    elif block.shape_variant == "editorial-panel":
        title_space_after = Pt(8)
        body_space_after = Pt(5)

    frame.clear()
    title_paragraph = frame.paragraphs[0]
    title_paragraph.text = block.title or ""
    title_paragraph.alignment = title_alignment
    title_paragraph.font.size = Pt(title_size)
    title_paragraph.font.bold = True
    title_paragraph.font.name = slide_spec.title_font_family
    title_paragraph.font.color.rgb = title_color
    title_paragraph.space_before = Pt(0)
    title_paragraph.space_after = title_space_after
    title_paragraph.line_spacing = title_line_spacing

    for line in block.text_lines:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.alignment = body_alignment
        paragraph.font.size = Pt(body_size)
        paragraph.font.name = slide_spec.body_font_family
        paragraph.font.color.rgb = body_color
        paragraph.space_before = Pt(0)
        paragraph.space_after = body_space_after
        paragraph.line_spacing = body_line_spacing


def _add_svg_block_to_ppt(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    block: SvgBlockSpec,
) -> None:
    left = _scale_x(presentation, block.x, slide_spec.width)
    top = _scale_y(presentation, block.y, slide_spec.height)
    width = _scale_x(presentation, block.width, slide_spec.width)
    height = _scale_y(presentation, block.height, slide_spec.height)

    _add_block_underlay(presentation, slide, slide_spec, block)
    shape = slide.shapes.add_shape(
        _shape_type(block.shape_variant),
        left,
        top,
        width,
        height,
    )
    _apply_block_shape_style(shape, block, slide_spec)
    _apply_block_text_style(shape.text_frame, block, slide_spec)
    _add_block_accents(presentation, slide, slide_spec, block)


def _apply_slide_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color, "#f7fbff")


def _citation_label(citation: Citation) -> str:
    label = citation.note or citation.page_label or citation.asset_id or "Reference"
    normalized = " ".join(label.split()).strip()
    return normalized[:48]


def _citation_layout_profile(slide_spec: SvgSlideSpec) -> dict[str, int | float | str]:
    profile: dict[str, int | float | str] = {
        "heading_x": slide_spec.width - 320,
        "heading_y": 88,
        "heading_width": 120,
        "heading_height": 24,
        "chip_width": 250,
        "chip_height": 30,
        "chip_gap": 8,
        "max_visible": 3,
        "chip_fill_transparency": 0.18,
        "chip_line_transparency": 0.48,
    }

    if slide_spec.layout_name == "cover-hero":
        profile.update(
            {
                "heading_x": slide_spec.width - 326,
                "heading_width": 136,
                "chip_width": 262,
                "chip_height": 32,
                "chip_gap": 10,
                "max_visible": 2,
            }
        )
    elif slide_spec.layout_name == "comparison-columns":
        profile.update(
            {
                "heading_x": slide_spec.width - 310,
                "heading_y": 104,
                "heading_width": 126,
                "chip_width": 236,
                "chip_height": 28,
                "max_visible": 2,
            }
        )
    elif slide_spec.layout_name == "timeline-ribbon":
        profile.update(
            {
                "heading_x": slide_spec.width - 296,
                "heading_y": 74,
                "heading_width": 118,
                "chip_width": 220,
                "chip_height": 26,
                "chip_gap": 6,
                "max_visible": 2,
                "chip_fill_transparency": 0.1,
                "chip_line_transparency": 0.62,
            }
        )
    elif slide_spec.layout_name == "process-ladder":
        profile.update(
            {
                "heading_x": slide_spec.width - 308,
                "heading_y": 162,
                "chip_width": 232,
                "max_visible": 2,
            }
        )
    elif slide_spec.layout_name in {"media-gallery", "workshop-board"}:
        profile.update(
            {
                "heading_x": slide_spec.width - 316,
                "heading_y": 96,
                "chip_width": 244,
                "max_visible": 2,
            }
        )
    elif slide_spec.layout_name == "assignment-brief":
        profile.update(
            {
                "heading_x": slide_spec.width - 314,
                "heading_y": 86,
                "chip_width": 238,
            }
        )
    elif slide_spec.layout_name == "recap-strip":
        profile.update(
            {
                "heading_x": slide_spec.width - 316,
                "heading_y": 82,
                "chip_width": 242,
                "max_visible": 2,
            }
        )

    return profile


def _citation_theme_style(theme_id: str | None, slide_spec: SvgSlideSpec) -> dict[str, str | float]:
    theme_style = dict(PPT_CITATION_THEME_STYLES.get(theme_id or "academy", PPT_CITATION_THEME_STYLES["academy"]))
    theme_style.setdefault("heading_fill", slide_spec.accent_color)
    theme_style.setdefault("heading_text", "#ffffff")
    theme_style.setdefault("heading_fill_transparency", 0.04)
    theme_style.setdefault("chip_fill", slide_spec.soft_color)
    theme_style.setdefault("chip_text", slide_spec.accent_color)
    theme_style.setdefault("chip_line", slide_spec.accent_color)
    theme_style.setdefault("chip_fill_transparency", 0.18)
    theme_style.setdefault("chip_line_transparency", 0.48)
    return theme_style


def _add_text_to_shape(
    shape,
    *,
    title: str,
    font_name: str,
    font_size: int,
    color: str | RGBColor,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bold: bool = True,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.alignment = align
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0


def _add_slide_citation_tags(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    citations: list[Citation],
    *,
    theme_id: str | None = None,
) -> None:
    labels: list[str] = []
    for citation in citations:
        label = _citation_label(citation)
        if label and label not in labels:
            labels.append(label)

    if not labels:
        return

    profile = _citation_layout_profile(slide_spec)
    theme_style = _citation_theme_style(theme_id, slide_spec)
    max_visible = int(profile["max_visible"])
    visible_labels = labels[:max_visible]
    remaining = len(labels) - len(visible_labels)
    if remaining > 0:
        visible_labels.append(f"+{remaining} more")

    heading_x = int(profile["heading_x"])
    heading_y = int(profile["heading_y"])
    heading_width = int(profile["heading_width"])
    heading_height = int(profile["heading_height"])
    chip_width = int(profile["chip_width"])
    chip_height = int(profile["chip_height"])
    chip_gap = int(profile["chip_gap"])

    if slide_spec.layout_name == "cover-hero":
        total_height = heading_height + 6 + (len(visible_labels) * chip_height) + (max(len(visible_labels) - 1, 0) * chip_gap)
        heading_y = max(slide_spec.height - total_height - 44, 86)

    heading = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _scale_x(presentation, heading_x, slide_spec.width),
        _scale_y(presentation, heading_y, slide_spec.height),
        _scale_x(presentation, heading_width, slide_spec.width),
        _scale_y(presentation, heading_height, slide_spec.height),
    )
    heading.fill.solid()
    heading.fill.fore_color.rgb = _rgb(str(theme_style["heading_fill"]), "#16324f")
    heading.fill.transparency = float(theme_style["heading_fill_transparency"])
    heading.line.fill.background()
    _add_text_to_shape(
        heading,
        title="References",
        font_name=slide_spec.title_font_family,
        font_size=11,
        color=str(theme_style["heading_text"]),
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    chip_y = heading_y + heading_height + 6
    for label in visible_labels:
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            _scale_x(presentation, heading_x, slide_spec.width),
            _scale_y(presentation, chip_y, slide_spec.height),
            _scale_x(presentation, chip_width, slide_spec.width),
            _scale_y(presentation, chip_height, slide_spec.height),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = _rgb(str(theme_style["chip_fill"]), "#d7e6f5")
        chip.fill.transparency = float(theme_style["chip_fill_transparency"])
        chip.line.color.rgb = _rgb(str(theme_style["chip_line"]), "#16324f")
        chip.line.transparency = float(theme_style["chip_line_transparency"])
        _add_text_to_shape(
            chip,
            title=label,
            font_name=slide_spec.body_font_family,
            font_size=9 if slide_spec.layout_name == "timeline-ribbon" else 10,
            color=str(theme_style["chip_text"]),
            align=PP_ALIGN.LEFT,
            bold=False,
        )
        chip_y += chip_height + chip_gap


def _add_svg_decorations(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
) -> None:
    accent = _rgb(slide_spec.accent_color, "#16324f")
    soft = _rgb(slide_spec.soft_color, "#d7e6f5")

    if slide_spec.layout_name == "cover-hero":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=48,
            y=150,
            width=760,
            height=338,
            fill_color=soft,
            transparency=0.72,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=72,
            y=134,
            width=460,
            height=18,
            fill_color=accent,
            transparency=0.84,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            _scale_x(presentation, 960, slide_spec.width),
            _scale_y(presentation, 24, slide_spec.height),
            _scale_x(presentation, 300, slide_spec.width),
            _scale_y(presentation, 300, slide_spec.height),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.fill.transparency = 0.9
        circle.line.fill.background()
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=104,
            y=82,
            width=220,
            height=10,
            fill_color=accent,
            transparency=0.7,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=1030,
            y=82,
            width=116,
            height=116,
            fill_color="#ffffff",
            transparency=0.74,
            shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=858,
            y=566,
            width=248,
            height=22,
            fill_color=accent,
            transparency=0.84,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
    elif slide_spec.layout_name == "process-ladder":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=184,
            y=168,
            width=14,
            height=410,
            fill_color=accent,
            transparency=0.84,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        for offset in (212, 344, 476):
            _add_overlay_shape(
                presentation,
                slide,
                slide_spec,
                x=162,
                y=offset,
                width=58,
                height=58,
                fill_color=soft,
                transparency=0.2,
                shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
            )
    elif slide_spec.layout_name == "media-gallery":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=706,
            y=152,
            width=438,
            height=318,
            fill_color=soft,
            transparency=0.55,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=736,
            y=172,
            width=378,
            height=12,
            fill_color=accent,
            transparency=0.72,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
    elif slide_spec.layout_name == "workshop-board":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=852,
            y=96,
            width=230,
            height=18,
            fill_color=accent,
            transparency=0.75,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=900,
            y=108,
            width=26,
            height=26,
            fill_color=accent,
            transparency=0.15,
            shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
        )
    elif slide_spec.layout_name == "assignment-brief":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=116,
            y=164,
            width=18,
            height=392,
            fill_color=accent,
            transparency=0.55,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=1034,
            y=108,
            width=96,
            height=96,
            fill_color=soft,
            transparency=0.18,
            shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
        )
    elif slide_spec.layout_name == "recap-strip":
        _add_overlay_shape(
            presentation,
            slide,
            slide_spec,
            x=112,
            y=522,
            width=1068,
            height=18,
            fill_color=accent,
            transparency=0.82,
            shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
    elif slide_spec.layout_name == "timeline-ribbon":
        ribbon = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            _scale_x(presentation, 120, slide_spec.width),
            _scale_y(presentation, 334, slide_spec.height),
            _scale_x(presentation, 1040, slide_spec.width),
            _scale_y(presentation, 10, slide_spec.height),
        )
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = accent
        ribbon.fill.transparency = 0.86
        ribbon.line.fill.background()
        for offset in (278, 532, 786):
            _add_overlay_shape(
                presentation,
                slide,
                slide_spec,
                x=offset,
                y=314,
                width=34,
                height=34,
                fill_color=soft,
                transparency=0.08,
                shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
            )
    else:
        blob = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            _scale_x(presentation, 980, slide_spec.width),
            _scale_y(presentation, 36, slide_spec.height),
            _scale_x(presentation, 220, slide_spec.width),
            _scale_y(presentation, 220, slide_spec.height),
        )
        blob.fill.solid()
        blob.fill.fore_color.rgb = soft
        blob.fill.transparency = 0.68
        blob.line.fill.background()


def _render_svg_slide_to_ppt(
    presentation: Presentation,
    slide,
    slide_spec: SvgSlideSpec,
    citations: list[Citation] | None = None,
    theme_id: str | None = None,
) -> None:
    _apply_slide_background(slide, slide_spec.background)
    _add_svg_decorations(presentation, slide, slide_spec)
    for block in slide_spec.blocks:
        _add_svg_block_to_ppt(presentation, slide, slide_spec, block)
    _add_slide_citation_tags(presentation, slide, slide_spec, citations or [], theme_id=theme_id)


def export_pptx_for_session(
    session: SessionState,
    store_namespace: str | None = None,
    top_k: int = 5,
    theme_id: str | None = None,
    font_preset: str | None = None,
) -> tuple[SessionState, ExportArtifact]:
    session = _ensure_svg_deck(
        session,
        store_namespace,
        top_k,
        theme_id=theme_id,
        font_preset=font_preset,
    )

    ensure_project_directories()
    lesson_title = session.teaching_spec.lesson_title or "lesson"
    export_path = build_export_path(session.session_id, f"{lesson_title}_slides", "pptx")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = f"{lesson_title} slide deck"
    presentation.core_properties.subject = session.teaching_spec.subject or "teaching-agent"
    presentation.core_properties.comments = "Generated by Teaching Agent"

    assert session.svg_deck is not None
    slide_lookup = {item.slide_number: item for item in session.slide_plan.slides}
    for slide_spec in session.svg_deck.slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        plan_item = slide_lookup.get(slide_spec.slide_number)
        _render_svg_slide_to_ppt(
            presentation,
            slide,
            slide_spec,
            citations=plan_item.citations if plan_item is not None else [],
            theme_id=session.svg_deck.theme_id,
        )

    presentation.save(export_path)

    artifact = ExportArtifact(
        filename=export_path.name,
        resource_type=ResourceType.PPTX,
        path=str(export_path),
        summary=_build_pptx_summary(session.teaching_spec, len(session.svg_deck.slides)),
    )
    return _finalize_export(session, artifact)
