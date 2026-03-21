from __future__ import annotations

import shutil
from pathlib import Path
from uuid import uuid4

from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor

from app.config import get_settings
from app.main import app
from app.models import Citation, InteractionMode, PlanningConfirmation, SessionState, SlidePlan, SlidePlanItem, SlideType, TeachingSpec
from app.services.exporter import export_docx_for_session, export_pptx_for_session
from app.services.storage import session_store


client = TestClient(app)


def _slide_texts(presentation: Presentation) -> list[str]:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return texts


def _find_shape_with_text(presentation: Presentation, needle: str):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame and needle in shape.text_frame.text:
                return shape
    raise AssertionError(f"Shape containing {needle!r} not found")


def _find_shape_with_text_in_slide(slide, needle: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    raise AssertionError(f"Shape containing {needle!r} not found in slide")


def _find_shape_with_text_in_slide_region(slide, needle: str, *, min_left: int, max_width: int | None = None, max_height: int | None = None):
    for shape in slide.shapes:
        if (
            getattr(shape, "has_text_frame", False)
            and shape.has_text_frame
            and needle in shape.text_frame.text
            and shape.left >= min_left
            and (max_width is None or shape.width <= max_width)
            and (max_height is None or shape.height <= max_height)
        ):
            return shape
    raise AssertionError(f"Shape containing {needle!r} not found in expected slide region")


def test_export_docx_for_session_writes_file() -> None:
    session = SessionState(
        title="Export Unit Test",
        teaching_spec=TeachingSpec(
            education_stage="middle-school",
            subject="history",
            lesson_title="工业革命",
            class_duration_minutes=45,
        ),
    )

    updated_session, artifact = export_docx_for_session(session)

    path = Path(artifact.path)
    assert path.exists()
    assert artifact.filename.endswith(".docx")
    assert updated_session.export_artifacts
    assert updated_session.stage.value == "export"

    document = Document(path)
    full_text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
    assert "工业革命" in full_text
    assert "逐页策划" in full_text

    shutil.rmtree(path.parent, ignore_errors=True)


def test_export_docx_api_and_download() -> None:
    settings = get_settings()
    source_dir = settings.knowledge_base_dir / f"_export_tests_{uuid4().hex}"
    source_dir.mkdir(parents=True, exist_ok=True)
    namespace = f"export_api_{uuid4().hex}"
    session_export_dir: Path | None = None

    try:
        (source_dir / "english.txt").write_text(
            "Environment Protection lessons can include vocabulary scaffolds, discussion prompts, and project tasks.",
            encoding="utf-8",
        )
        ingest_response = client.post(
            "/api/kb/ingest",
            json={
                "source_dir": str(source_dir),
                "reset": True,
                "store_namespace": namespace,
            },
        )
        assert ingest_response.status_code == 200

        chat_response = client.post(
            "/api/chat/messages",
            json={
                "title": "Export Demo",
                "content": '我想做一节高中英语"Environment Protection"课程，40分钟，教学目标：理解环境问题并组织讨论。加入讨论和项目任务。',
            },
        )
        assert chat_response.status_code == 200
        session_id = chat_response.json()["session_id"]
        confirm_response = client.post(
            "/api/planner/confirmation/confirm",
            json={"session_id": session_id, "note": "约束完整后导出正式稿"},
        )
        assert confirm_response.status_code == 200

        export_response = client.post(
            "/api/export/docx",
            json={
                "session_id": session_id,
                "store_namespace": namespace,
                "top_k": 3,
            },
        )
        assert export_response.status_code == 200
        payload = export_response.json()
        assert payload["artifact"]["filename"].endswith(".docx")
        assert payload["download_url"].startswith(f"/api/export/files/{session_id}/")
        assert payload["session"]["export_artifacts"]

        artifact_path = Path(payload["artifact"]["path"])
        session_export_dir = artifact_path.parent
        assert artifact_path.exists()

        download_response = client.get(payload["download_url"])
        assert download_response.status_code == 200
        assert (
            download_response.headers["content-type"]
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
    finally:
        shutil.rmtree(source_dir, ignore_errors=True)
        shutil.rmtree(settings.vector_store_dir / namespace, ignore_errors=True)
        if session_export_dir is not None:
            shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_for_session_writes_file() -> None:
    session = SessionState(
        title="PPTX Export Unit Test",
        teaching_spec=TeachingSpec(
            education_stage="middle-school",
            subject="math",
            lesson_title="一次函数",
            class_duration_minutes=50,
        ),
    )

    updated_session, artifact = export_pptx_for_session(session)

    path = Path(artifact.path)
    assert path.exists()
    assert artifact.filename.endswith(".pptx")
    assert updated_session.export_artifacts
    assert updated_session.stage.value == "export"

    presentation = Presentation(path)
    assert len(presentation.slides) >= 1
    text_blob = "\n".join(_slide_texts(presentation))
    assert "一次函数" in text_blob
    assert "Knowledge Core" in text_blob or "Class Snapshot" in text_blob or "Big Takeaway" in text_blob

    shutil.rmtree(path.parent, ignore_errors=True)


def test_export_pptx_api_and_download() -> None:
    settings = get_settings()
    source_dir = settings.knowledge_base_dir / f"_pptx_export_tests_{uuid4().hex}"
    source_dir.mkdir(parents=True, exist_ok=True)
    namespace = f"export_pptx_api_{uuid4().hex}"
    session_export_dir: Path | None = None

    try:
        (source_dir / "history.txt").write_text(
            "Industrial Revolution lessons benefit from timeline evidence, source analysis, and discussion prompts.",
            encoding="utf-8",
        )
        ingest_response = client.post(
            "/api/kb/ingest",
            json={
                "source_dir": str(source_dir),
                "reset": True,
                "store_namespace": namespace,
            },
        )
        assert ingest_response.status_code == 200

        chat_response = client.post(
            "/api/chat/messages",
            json={
                "title": "PPTX Export Demo",
                "content": "我想做一节初中历史《工业革命》课程，45分钟，教学目标：理解蒸汽机与工厂制度的关系。加入材料分析和讨论。",
            },
        )
        assert chat_response.status_code == 200
        session_id = chat_response.json()["session_id"]
        confirm_response = client.post(
            "/api/planner/confirmation/confirm",
            json={"session_id": session_id, "note": "确认后导出正式课件"},
        )
        assert confirm_response.status_code == 200

        export_response = client.post(
            "/api/export/pptx",
            json={
                "session_id": session_id,
                "store_namespace": namespace,
                "top_k": 3,
            },
        )
        assert export_response.status_code == 200
        payload = export_response.json()
        assert payload["artifact"]["filename"].endswith(".pptx")
        assert payload["download_url"].startswith(f"/api/export/files/{session_id}/")
        assert payload["session"]["export_artifacts"]

        artifact_path = Path(payload["artifact"]["path"])
        session_export_dir = artifact_path.parent
        assert artifact_path.exists()
        presentation = Presentation(artifact_path)
        assert len(presentation.slides) == len(payload["session"]["svg_deck"]["slides"])
        assert "工业革命" in "\n".join(_slide_texts(presentation))

        download_response = client.get(payload["download_url"])
        assert download_response.status_code == 200
        assert (
            download_response.headers["content-type"]
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    finally:
        shutil.rmtree(source_dir, ignore_errors=True)
        shutil.rmtree(settings.vector_store_dir / namespace, ignore_errors=True)
        if session_export_dir is not None:
            shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_api_blocks_unconfirmed_session() -> None:
    chat_response = client.post(
        "/api/chat/messages",
        json={
            "title": "Blocked Export Demo",
            "content": '我想做一节高中英语"Environment Protection"课程，40分钟，加入讨论和项目任务。',
        },
    )
    assert chat_response.status_code == 200
    session_id = chat_response.json()["session_id"]

    export_response = client.post(
        "/api/export/docx",
        json={"session_id": session_id, "top_k": 3},
    )

    assert export_response.status_code == 409
    assert "关键约束尚未确认" in export_response.json()["detail"]


def test_export_api_blocks_high_risk_content_even_after_confirmation() -> None:
    session = SessionState(
        title="Blocked High Risk Export",
        teaching_spec=TeachingSpec(
            education_stage="high-school",
            subject="english",
            lesson_title="Environment Protection",
            class_duration_minutes=40,
            learning_objectives=[{"description": "围绕环境议题组织阅读与讨论"}],
        ),
        slide_plan=SlidePlan(
            title="Environment Protection",
            slides=[
                SlidePlanItem(
                    slide_number=1,
                    title="Environment Protection",
                    slide_type=SlideType.COVER,
                    goal="建立课堂主题",
                    key_points=["topic framing"],
                ),
                SlidePlanItem(
                    slide_number=2,
                    title="Content Expansion",
                    slide_type=SlideType.CONCEPT,
                    goal="展开正文内容",
                    key_points=["解释环境保护的多重社会影响"],
                    speaker_notes=["继续展开正文但没有证据来源"],
                    citations=[],
                ),
            ],
        ),
        planning_confirmation=PlanningConfirmation(confirmed=True, summary="约束已确认"),
    )
    session_store.save(session)

    try:
        export_response = client.post(
            "/api/export/pptx",
            json={"session_id": session.session_id, "top_k": 3},
        )
        assert export_response.status_code == 409
        assert "已拦截正式导出" in export_response.json()["detail"]
    finally:
        settings = get_settings()
        shutil.rmtree(settings.workspaces_dir / session.session_id, ignore_errors=True)
        shutil.rmtree(settings.exports_dir / session.session_id, ignore_errors=True)
        shutil.rmtree(settings.raw_data_dir / session.session_id, ignore_errors=True)
        shutil.rmtree(settings.parsed_data_dir / session.session_id, ignore_errors=True)


def test_export_pptx_preserves_svg_theme_preferences_after_plan_changes() -> None:
    chat_response = client.post(
        "/api/chat/messages",
        json={
            "title": "Theme Persistence Demo",
            "content": "我想做一节初中历史《工业革命》课程，45分钟，教学目标：理解蒸汽机与工厂制度的关系。加入材料分析和讨论。",
        },
    )
    assert chat_response.status_code == 200
    session_id = chat_response.json()["session_id"]

    session_export_dir: Path | None = None

    try:
        svg_response = client.post(
            "/api/svg/deck",
            json={
                "session_id": session_id,
                "top_k": 5,
                "theme_id": "studio",
                "font_preset": "modern",
            },
        )
        assert svg_response.status_code == 200
        svg_payload = svg_response.json()
        assert svg_payload["session"]["svg_theme_id"] == "studio"
        assert svg_payload["session"]["svg_font_preset"] == "modern"

        update_response = client.post(
            "/api/planner/slide-plan/update",
            json={
                "session_id": session_id,
                "slide_number": 2,
                "title": "工业革命核心线索",
                "goal": "按技术、制度和社会变化重新组织本页内容",
                "revision_note": "invalidate svg deck",
            },
        )
        assert update_response.status_code == 200
        updated_payload = update_response.json()
        assert updated_payload["session"]["svg_deck"] is None
        assert updated_payload["session"]["svg_theme_id"] == "studio"
        assert updated_payload["session"]["svg_font_preset"] == "modern"

        confirm_response = client.post(
            "/api/planner/confirmation/confirm",
            json={"session_id": session_id, "note": "确认约束后继续导出"},
        )
        assert confirm_response.status_code == 200

        export_response = client.post(
            "/api/export/pptx",
            json={
                "session_id": session_id,
                "top_k": 5,
            },
        )
        assert export_response.status_code == 200
        payload = export_response.json()
        assert payload["session"]["svg_theme_id"] == "studio"
        assert payload["session"]["svg_font_preset"] == "modern"
        assert payload["session"]["svg_deck"]["theme_id"] == "studio"
        assert payload["session"]["svg_deck"]["font_preset"] == "modern"
        assert payload["session"]["svg_deck"]["slides"][0]["title_font_family"] == "Bahnschrift"

        artifact_path = Path(payload["artifact"]["path"])
        session_export_dir = artifact_path.parent
        assert artifact_path.exists()
    finally:
        if session_export_dir is not None:
            shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_high_fidelity_layouts_add_extra_shapes() -> None:
    session = SessionState(
        title="High Fidelity PPTX Export",
        teaching_spec=TeachingSpec(
            education_stage="high-school",
            subject="english",
            lesson_title="Environment Protection",
            class_duration_minutes=40,
        ),
        slide_plan=SlidePlan(
            title="Environment Protection",
            slides=[
                SlidePlanItem(
                    slide_number=1,
                    title="Lesson Flow",
                    slide_type=SlideType.PROCESS,
                    goal="把课堂流程拆成清晰的活动梯度",
                    interaction_mode=InteractionMode.DISCUSSION,
                    key_points=["Warm-up", "Input", "Practice"],
                    visual_brief=["step ladder"],
                    speaker_notes=["guide transitions"],
                ),
                SlidePlanItem(
                    slide_number=2,
                    title="Visual Sources",
                    slide_type=SlideType.MEDIA,
                    goal="组合图像、数据和提示语",
                    interaction_mode=InteractionMode.DISCUSSION,
                    key_points=["photo prompt", "data card", "caption"],
                    visual_brief=["gallery wall"],
                    speaker_notes=["compare media evidence"],
                ),
                SlidePlanItem(
                    slide_number=3,
                    title="Workshop Sprint",
                    slide_type=SlideType.ACTIVITY,
                    goal="组织小组任务与分工",
                    interaction_mode=InteractionMode.PROJECT,
                    key_points=["task board", "roles", "deliverable"],
                    visual_brief=["studio board"],
                    speaker_notes=["check group pacing"],
                ),
                SlidePlanItem(
                    slide_number=4,
                    title="After Class",
                    slide_type=SlideType.ASSIGNMENT,
                    goal="布置课后延伸任务",
                    interaction_mode=InteractionMode.EXERCISE,
                    key_points=["reflection", "mini project", "submission"],
                    visual_brief=["brief panel"],
                    speaker_notes=["clarify evaluation"],
                ),
                SlidePlanItem(
                    slide_number=5,
                    title="History of Change",
                    slide_type=SlideType.TIMELINE,
                    goal="回顾关键阶段与转折",
                    interaction_mode=InteractionMode.DEBATE,
                    key_points=["past", "present", "future"],
                    visual_brief=["timeline ribbon"],
                    speaker_notes=["connect milestones"],
                ),
            ],
        ),
    )

    updated_session, artifact = export_pptx_for_session(session, theme_id="studio", font_preset="modern")

    path = Path(artifact.path)
    session_export_dir = path.parent

    try:
        presentation = Presentation(path)
        text_blob = "\n".join(_slide_texts(presentation))
        assert "Method Frame" in text_blob
        assert "Visual Stage" in text_blob
        assert "Studio Task" in text_blob
        assert "After Class" in text_blob
        assert "Timeline Lens" in text_blob

        assert updated_session.svg_deck is not None
        assert any(
            len(presentation.slides[index].shapes) > len(slide_spec.blocks) + 1
            for index, slide_spec in enumerate(updated_session.svg_deck.slides)
        )
    finally:
        shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_cover_typography_uses_hierarchy_and_spacing() -> None:
    session = SessionState(
        title="Cover Typography Demo",
        teaching_spec=TeachingSpec(
            education_stage="middle-school",
            subject="science",
            lesson_title="生态系统",
            class_duration_minutes=45,
        ),
    )

    updated_session, artifact = export_pptx_for_session(session, theme_id="academy", font_preset="reading")
    path = Path(artifact.path)
    session_export_dir = path.parent

    try:
        presentation = Presentation(path)
        hero_shape = _find_shape_with_text(presentation, "生态系统")
        title_paragraph = hero_shape.text_frame.paragraphs[0]
        body_paragraphs = [paragraph for paragraph in hero_shape.text_frame.paragraphs[1:] if paragraph.text]

        assert title_paragraph.font.size is not None
        assert title_paragraph.font.size.pt >= 30
        assert title_paragraph.space_after is not None
        assert body_paragraphs
        assert all(paragraph.space_after is not None for paragraph in body_paragraphs)

        assert updated_session.svg_deck is not None
        assert len(presentation.slides[0].shapes) > len(updated_session.svg_deck.slides[0].blocks) + 3
    finally:
        shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_renders_page_level_citation_tags() -> None:
    session = SessionState(
        title="Citation Tag Demo",
        teaching_spec=TeachingSpec(
            education_stage="middle-school",
            subject="history",
            lesson_title="工业革命",
            class_duration_minutes=45,
        ),
        slide_plan=SlidePlan(
            title="工业革命",
            slides=[
                SlidePlanItem(
                    slide_number=1,
                    title="工业革命导入",
                    slide_type=SlideType.COVER,
                    goal="建立主题背景和问题意识",
                    key_points=["技术变革", "社会影响"],
                    speaker_notes=["提出导入问题"],
                    citations=[
                        Citation(asset_id="hist-1", page_label="教材第12页", note="教材第12页"),
                        Citation(asset_id="hist-2", page_label="史料图表A", note="史料图表A"),
                    ],
                ),
                SlidePlanItem(
                    slide_number=2,
                    title="影响分析",
                    slide_type=SlideType.COMPARISON,
                    goal="比较工业革命前后的社会变化",
                    key_points=["生产方式", "城市化", "劳动结构"],
                    speaker_notes=["结合材料引导比较"],
                    citations=[
                        Citation(asset_id="hist-3", page_label="资料卡B", note="资料卡B"),
                    ],
                ),
            ],
        ),
    )

    updated_session, artifact = export_pptx_for_session(session, theme_id="academy", font_preset="reading")
    path = Path(artifact.path)
    session_export_dir = path.parent

    try:
        presentation = Presentation(path)
        text_blob = "\n".join(_slide_texts(presentation))
        assert "References" in text_blob
        assert "教材第12页" in text_blob
        assert "史料图表A" in text_blob
        assert "资料卡B" in text_blob

        assert updated_session.svg_deck is not None
        first_slide = presentation.slides[0]
        assert len(first_slide.shapes) > len(updated_session.svg_deck.slides[0].blocks) + 4
    finally:
        shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_citation_tags_follow_layout_specific_positions_and_density() -> None:
    session = SessionState(
        title="Citation Layout Demo",
        teaching_spec=TeachingSpec(
            education_stage="middle-school",
            subject="history",
            lesson_title="工业革命",
            class_duration_minutes=45,
        ),
        slide_plan=SlidePlan(
            title="工业革命",
            slides=[
                SlidePlanItem(
                    slide_number=1,
                    title="工业革命导入",
                    slide_type=SlideType.COVER,
                    goal="建立主题背景和问题意识",
                    key_points=["技术变革", "社会影响"],
                    citations=[
                        Citation(asset_id="hist-1", note="教材第12页"),
                        Citation(asset_id="hist-2", note="史料图表A"),
                        Citation(asset_id="hist-3", note="课堂观察单"),
                    ],
                ),
                SlidePlanItem(
                    slide_number=2,
                    title="影响比较",
                    slide_type=SlideType.COMPARISON,
                    goal="比较工业革命前后的变化",
                    key_points=["生产方式", "城市化", "劳动结构"],
                    citations=[
                        Citation(asset_id="hist-4", note="资料卡B"),
                        Citation(asset_id="hist-5", note="数据表C"),
                    ],
                ),
                SlidePlanItem(
                    slide_number=3,
                    title="关键阶段",
                    slide_type=SlideType.TIMELINE,
                    goal="梳理关键时间节点",
                    key_points=["起点", "扩散", "影响"],
                    citations=[
                        Citation(asset_id="hist-6", note="时间轴资料1"),
                        Citation(asset_id="hist-7", note="时间轴资料2"),
                        Citation(asset_id="hist-8", note="时间轴资料3"),
                        Citation(asset_id="hist-9", note="时间轴资料4"),
                    ],
                ),
            ],
        ),
    )

    _, artifact = export_pptx_for_session(session, theme_id="academy", font_preset="reading")
    path = Path(artifact.path)
    session_export_dir = path.parent

    try:
        presentation = Presentation(path)
        cover_ref = _find_shape_with_text_in_slide(presentation.slides[0], "References")
        comparison_ref = _find_shape_with_text_in_slide(presentation.slides[1], "References")
        timeline_ref = _find_shape_with_text_in_slide(presentation.slides[2], "References")

        assert cover_ref.top > int(presentation.slide_height * 0.6)
        assert comparison_ref.top < int(presentation.slide_height * 0.3)
        assert timeline_ref.top < int(presentation.slide_height * 0.22)

        text_blob = "\n".join(_slide_texts(presentation))
        assert "+1 more" in text_blob
        assert "+2 more" in text_blob
    finally:
        shutil.rmtree(session_export_dir, ignore_errors=True)


def test_export_pptx_citation_tags_follow_theme_styles() -> None:
    session = SessionState(
        title="Citation Theme Demo",
        teaching_spec=TeachingSpec(
            education_stage="high-school",
            subject="english",
            lesson_title="Environment Protection",
            class_duration_minutes=40,
        ),
        slide_plan=SlidePlan(
            title="Environment Protection",
            slides=[
                SlidePlanItem(
                    slide_number=1,
                    title="Environment Protection",
                    slide_type=SlideType.COVER,
                    goal="建立环境议题的课堂入口",
                    key_points=["topic framing", "visual prompt"],
                    citations=[
                        Citation(asset_id="eng-1", note="Workbook p.14"),
                        Citation(asset_id="eng-2", note="Photo Set A"),
                    ],
                )
            ],
        ),
    )

    studio_session, studio_artifact = export_pptx_for_session(session.model_copy(deep=True), theme_id="studio", font_preset="modern")
    briefing_session, briefing_artifact = export_pptx_for_session(session.model_copy(deep=True), theme_id="briefing", font_preset="modern")

    studio_path = Path(studio_artifact.path)
    briefing_path = Path(briefing_artifact.path)

    try:
        studio_presentation = Presentation(studio_path)
        briefing_presentation = Presentation(briefing_path)

        studio_ref = _find_shape_with_text_in_slide(studio_presentation.slides[0], "References")
        briefing_ref = _find_shape_with_text_in_slide(briefing_presentation.slides[0], "References")
        min_left = int(studio_presentation.slide_width * 0.68)
        max_width = int(studio_presentation.slide_width * 0.24)
        max_height = int(studio_presentation.slide_height * 0.07)
        studio_chip = _find_shape_with_text_in_slide_region(
            studio_presentation.slides[0],
            "Workbook p.14",
            min_left=min_left,
            max_width=max_width,
            max_height=max_height,
        )
        briefing_chip = _find_shape_with_text_in_slide_region(
            briefing_presentation.slides[0],
            "Workbook p.14",
            min_left=min_left,
            max_width=max_width,
            max_height=max_height,
        )

        assert studio_session.svg_deck is not None
        assert briefing_session.svg_deck is not None
        assert studio_session.svg_deck.theme_id == "studio"
        assert briefing_session.svg_deck.theme_id == "briefing"

        assert studio_ref.fill.fore_color.rgb == RGBColor(0x8A, 0x2C, 0x0D)
        assert briefing_ref.fill.fore_color.rgb == RGBColor(0x0F, 0x17, 0x2A)
        assert studio_chip.fill.fore_color.rgb == RGBColor(0xFD, 0xE6, 0xD8)
        assert briefing_chip.fill.fore_color.rgb == RGBColor(0xE2, 0xE8, 0xF0)
        assert studio_chip.line.color.rgb == RGBColor(0x8A, 0x2C, 0x0D)
        assert briefing_chip.line.color.rgb == RGBColor(0x33, 0x41, 0x55)
    finally:
        shutil.rmtree(studio_path.parent, ignore_errors=True)
        shutil.rmtree(briefing_path.parent, ignore_errors=True)
