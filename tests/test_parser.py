from __future__ import annotations

from pathlib import Path
import shutil
from uuid import uuid4
import wave

import fitz
from PIL import Image
from docx import Document
from pptx import Presentation

from app.config import get_settings
from app.models import ResourceType
from app.services.parser import parse_file


def _make_test_dir() -> Path:
    root = get_settings().raw_data_dir / "_parser_tests"
    root.mkdir(parents=True, exist_ok=True)
    path = root / uuid4().hex
    path.mkdir(parents=True, exist_ok=True)
    return path


def test_parse_pdf_extracts_text() -> None:
    tmp_path = _make_test_dir()
    pdf_path = tmp_path / "lesson.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Industrial Revolution reshaped production and society.")
    doc.save(pdf_path)
    doc.close()

    try:
        parsed = parse_file(pdf_path)
        assert parsed.resource_type == ResourceType.PDF
        assert "Industrial Revolution" in parsed.extracted_text
        assert parsed.page_count == 1
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)


def test_parse_image_extracts_metadata() -> None:
    tmp_path = _make_test_dir()
    image_path = tmp_path / "diagram.png"
    image = Image.new("RGB", (320, 180), color="white")
    image.save(image_path)

    try:
        parsed = parse_file(image_path)
        assert parsed.resource_type == ResourceType.IMAGE
        assert parsed.width == 320
        assert parsed.height == 180
        assert "OCR not implemented yet" in parsed.warnings[0]
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)


def test_parse_audio_extracts_duration() -> None:
    tmp_path = _make_test_dir()
    audio_path = tmp_path / "sample.wav"
    with wave.open(str(audio_path), "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(16000)
        wav_file.writeframes(b"\x00\x00" * 16000)

    try:
        parsed = parse_file(audio_path)
        assert parsed.resource_type == ResourceType.AUDIO
        assert parsed.duration_seconds == 1.0
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)


def test_parse_docx_extracts_text() -> None:
    tmp_path = _make_test_dir()
    docx_path = tmp_path / "lesson.docx"
    document = Document()
    document.add_heading("环境保护", level=1)
    document.add_paragraph("本节课聚焦节能减排与可持续发展。")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "重点"
    table.rows[0].cells[1].text = "绿色生活方式"
    document.save(docx_path)

    try:
        parsed = parse_file(docx_path)
        assert parsed.resource_type == ResourceType.DOCX
        assert "环境保护" in parsed.extracted_text
        assert "绿色生活方式" in parsed.extracted_text
        assert parsed.metadata["paragraph_count"] >= 2
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)


def test_parse_pptx_extracts_slide_text() -> None:
    tmp_path = _make_test_dir()
    pptx_path = tmp_path / "lesson.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "一次函数"
    slide.placeholders[1].text = "一次函数的图像是一条直线。"
    presentation.save(pptx_path)

    try:
        parsed = parse_file(pptx_path)
        assert parsed.resource_type == ResourceType.PPTX
        assert "一次函数" in parsed.extracted_text
        assert parsed.page_count == 1
        assert parsed.metadata["slide_count"] == 1
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)
